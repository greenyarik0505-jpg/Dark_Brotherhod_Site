import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Function to add a styled title slide
def add_title_slide(title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 31, 53) # Dark blue/navy
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    subtitle.text = subtitle_text
    
    # Style title
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 215, 0) # Gold
    title.text_frame.paragraphs[0].font.bold = True
    
    # Style subtitle
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Function to add bullet slide
def add_bullet_slide(title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 31, 53)
    
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 215, 0)
    title.text_frame.paragraphs[0].font.bold = True
    
    tf = body.text_frame
    tf.text = bullet_points[0]
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.font.color.rgb = RGBColor(255, 255, 255)

# Slide 1
add_title_slide("Священная Империя", "Презентация официального сайта клуба\nСовременное решение для топового клуба в Brawl Stars")

# Slide 2
add_bullet_slide(
    "Зачем нам свой сайт?",
    [
        "Престиж и Статус: Уникальное лицо клуба в интернете, выделяющее нас среди тысяч других.",
        "Единая база: Вся важная информация, новости и правила собраны в одном месте.",
        "Мотивация участников: Зал славы, достижения и галерея вдохновляют игроков.",
        "Рекрутинг: Красивая страница 'Вступить' с требованиями помогает привлекать сильных игроков."
    ]
)

# Slide 3
add_bullet_slide(
    "Что умеет наш сайт?",
    [
        "Новости: Актуальные анонсы, результаты Лиги Клубов и объявления.",
        "Участники: Умный список с фильтрацией по ролям и сортировкой.",
        "События: Интерактивный таймлайн турниров и мегакопилок.",
        "Зал славы: Галерея скриншотов и рекорды клуба."
    ]
)

# Slide 4
add_bullet_slide(
    "Панель Управления (Админка)",
    [
        "Скрытая админка, доступная только по секретному паролю.",
        "Мгновенное редактирование, добавление и удаление данных.",
        "Удобный интерфейс — справится любой без знания кода.",
        "Все изменения мгновенно сохраняются в облаке (Google Firebase)."
    ]
)

# Slide 5
add_bullet_slide(
    "Империя ждет новых побед!",
    [
        "Сайт полностью готов к работе и адаптивен для телефонов.",
        "Высокая скорость загрузки и современный дизайн.",
        "Уже сейчас можно отправлять ссылку участникам клуба!"
    ]
)

prs.save("BrawlStars_Club_Presentation.pptx")
print("Presentation generated successfully.")
